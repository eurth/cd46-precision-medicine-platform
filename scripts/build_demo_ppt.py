#!/usr/bin/env python3
"""Generate OncoBridge demo PowerPoint. Run: python scripts/build_demo_ppt.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "OncoBridge_Demo_Presentation.pptx"

# Clinical Slate palette
NAVY = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
SLATE = RGBColor(0x64, 0x74, 0x8B)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

TARGETS = ("FOLH1 (PSMA)", "FAP", "SSTR2", "GRPR", "CD46")


def _box(slide, left, top, width, height, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb


def title_slide(prs, title, subtitle_lines=()):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    _box(slide, 0.6, 1.85, 8.8, 1.0, title, size=40, bold=True, color=WHITE)
    y = 3.05
    for line in subtitle_lines:
        _box(slide, 0.6, y, 8.8, 0.45, line, size=19, color=MUTED)
        y += 0.42
    _box(slide, 0.6, 6.85, 8.0, 0.35, "oncobridge.eurthtech.com  ·  Research use only", size=12, color=SLATE)


def content_slide(prs, title, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    _box(slide, 0.5, 0.35, 9.0, 0.7, title, size=28, bold=True, color=NAVY)
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(1.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    y = 1.25
    for i, b in enumerate(bullets):
        if not b:
            y += 0.25
            continue
        prefix = "• " if not b.startswith("  ") else "   "
        text = b.strip() if not b.startswith("  ") else b.strip()
        if not text.startswith("•"):
            text = prefix + text
        _box(slide, 0.55, y, 8.9, 0.55, text, size=16, color=NAVY if not b.startswith("  ") else SLATE)
        y += 0.52 if len(text) < 72 else 0.68
    if footer:
        _box(slide, 0.5, 6.9, 9.0, 0.4, footer, size=11, color=SLATE)


def target_grid_slide(prs):
    """Five targets — equal visual weight."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    _box(slide, 0.5, 0.35, 9.0, 0.7, "Five theranostic targets — equal platform scope", size=28, bold=True, color=NAVY)
    specs = [
        ("FOLH1 (PSMA)", "Prostate radioligand benchmark · Pluvicto® class"),
        ("FAP", "Stromal / CAF target · FAPI PET & emerging RLT"),
        ("SSTR2", "Neuroendocrine GPCR · somatostatin analog RLT"),
        ("GRPR", "Solid-tumour GPCR · bombesin-analogue radioligands"),
        ("CD46", "Pan-cancer surface antigen · complement pathway"),
    ]
    y = 1.35
    for sym, desc in specs:
        _box(slide, 0.55, y, 2.0, 0.35, sym, size=17, bold=True, color=BLUE)
        _box(slide, 2.6, y, 6.8, 0.45, desc, size=15, color=NAVY)
        y += 0.95
    _box(slide, 0.55, 6.75, 8.9, 0.4,
          "Switch any module via the target bar — same datasets, same graph schema, same sixteen modules.",
          size=12, color=SLATE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        "OncoBridge Intelligence",
        [
            "Multi-target theranostics research workbench",
            "FOLH1 (PSMA)  ·  FAP  ·  SSTR2  ·  GRPR  ·  CD46",
            "16 modules  ·  Knowledge graph  ·  Retrieval-augmented AI",
        ],
    )

    content_slide(prs, "The problem", [
        "Radioligand programs need expression, safety, patients, trials, and literature — in one place",
        "Public data lives in TCGA, GENIE, HPA, ChEMBL, ClinicalTrials.gov — disconnected formats",
        "Generic AI generates plausible oncology text without verifiable numbers",
        "Teams need one workbench across multiple surface targets — not one gene at a time",
    ])

    content_slide(prs, "What OncoBridge is", [
        "Open research platform — not a hospital system or medical device",
        "Five theranostic targets on one shared architecture",
        "16 modules across 9 research dimensions",
        "Neo4j knowledge graph (~3,586 nodes) — genes, drugs, trials, publications linked",
        "Target bar switches the active gene — charts, KG queries, and AI follow your selection",
    ], footer="Data freeze: 2026-07-28-phase4-five-targets")

    target_grid_slide(prs)

    content_slide(prs, "Nine research dimensions", [
        "Home — Platform orientation & module map",
        "Target / Cancer — Expression Atlas, Compare Targets",
        "Biomarkers — Multi-marker clinical decision views",
        "Proteins — PPI network, Diagnostics & imaging",
        "Patients — GENIE cohorts, Eligibility scoring",
        "Survival — Cox hazard ratios, forest plots",
        "Drugs / Safety — Pipeline, Dosimetry",
        "Graph / Ask — KG visualizer, Query Explorer, AI Assistant",
        "Strategy — End-to-end clinical development narrative",
    ])

    content_slide(prs, "Datasets included", [
        "TCGA / UCSC Xena — pan-cancer RNA & survival",
        "AACR GENIE — 271,837 real-world sequenced tumours",
        "Human Protein Atlas · GTEx · DepMap",
        "ClinicalTrials.gov · ChEMBL · Open Targets · PubMed · STRING · UniProt",
        "Neo4j AuraDB — OncoBridge knowledge graph",
        "Same provenance layer for every target in the registry",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    _box(slide, 0.5, 0.35, 9.0, 0.7, "ChatGPT vs OncoBridge", size=28, bold=True, color=NAVY)
    _box(slide, 0.5, 1.2, 4.2, 0.4, "Generic ChatGPT", size=16, bold=True, color=BLUE)
    for i, t in enumerate([
        "Generates plausible paragraphs",
        "May invent NCT IDs and hazard ratios",
        "No link to your TCGA extracts",
        "Cannot query your knowledge graph",
    ]):
        _box(slide, 0.5, 1.65 + i * 0.48, 4.3, 0.5, "• " + t, size=14, color=NAVY)
    _box(slide, 5.0, 1.2, 4.2, 0.4, "OncoBridge", size=16, bold=True, color=BLUE)
    for i, t in enumerate([
        "Retrieves ranked TCGA rows per active target",
        "Returns verifiable NCT trial IDs",
        "Pre-computed Cox HR from your CSVs",
        "Live Cypher — any of the five targets",
    ]):
        _box(slide, 5.0, 1.65 + i * 0.48, 4.3, 0.5, "• " + t, size=14, color=NAVY)

    content_slide(prs, "Live demo flow (20 min)", [
        "1. Platform Overview — five targets, module map, KG scale",
        "2. Compare Targets — all five genes on one chart (lead here)",
        "3. Expression Atlas — switch target bar: FOLH1 → FAP → SSTR2",
        "4. KG Query Explorer — Drugs template per target (rotate)",
        "5. Research Assistant — preset for whichever target is selected",
        "6. Patient Selection — GENIE cohort scale",
        "7. Clinical Strategy Engine — end-to-end narrative",
    ], footer="Script: reports/DEMO_RUNBOOK.md")

    content_slide(prs, "KG Query Explorer — rotate targets", [
        "Navigate: Graph → KG Query Explorer",
        "Select target in the bar (e.g. FOLH1), then:",
        "Template: 💊 Drugs: Agents targeting FOLH1?",
        "Run Query → drug name, type, max phase",
        "",
        "Repeat for FAP and SSTR2 — same templates, different symbol",
        "Say: “One graph schema — switch the gene, rerun the query.”",
    ])

    content_slide(prs, "KG Query Explorer — trials & expression", [
        "FOLH1: 🧪 Clinical Trials template → NCT IDs, phase, status",
        "FAP: 📊 Cell lines depending on FAP → DepMap CRISPR scores",
        "SSTR2: 🎯 Expression — highest cancers for SSTR2",
        "GRPR: 📚 Publications linked to GRPR",
        "CD46: 📈 Survival — High = worse prognosis",
    ])

    content_slide(prs, "Research Assistant — per active target", [
        "Select target in bar first — presets update to that gene",
        "",
        "FOLH1: “Summarise FOLH1 expression across TCGA cancer types with hazard ratios.”",
        "FAP: “What DepMap evidence supports FAP as a cancer dependency?”",
        "SSTR2: “What is the current SSTR2-targeted drug pipeline and clinical trials?”",
        "",
        "Say: “Retrieval-augmented — CSVs and graph for the active target.”",
    ], footer="Full list: reports/DEMO_QUESTIONS.md")

    content_slide(prs, "Acronym cheat sheet", [
        "TCGA — US cancer genomics atlas",
        "RLT — radioligand therapy",
        "PSMA / FOLH1 — prostate surface target",
        "FAP — fibroblast activation protein (stroma)",
        "SSTR2 — somatostatin receptor (neuroendocrine)",
        "GRPR — gastrin-releasing peptide receptor",
        "GENIE — AACR real-world sequencing registry",
        "NCT — clinical trial registry number",
    ])

    content_slide(prs, "Platform positioning", [
        "One workbench — five theranostic surface targets",
        "Same modules, datasets, and knowledge graph for each registered gene",
        "Target bar is the primary switch — audience should see you rotate it",
        "Research use only — not clinical advice, not a medical device",
    ])

    title_slide(
        prs,
        "Thank you",
        [
            "OncoBridge Intelligence",
            "oncobridge.eurthtech.com",
            "Questions?",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
    assert OUT.exists() and OUT.stat().st_size > 5000
